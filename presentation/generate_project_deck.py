#!/usr/bin/env python3
"""Generate the CFD9 project-history presentation."""

from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory

from pygments import highlight
from pygments.formatters import ImageFormatter
from pygments.lexers import FortranLexer, PythonLexer
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentation" / "CFD9_UCNS3D_project_history.pptx"

NAVY = RGBColor(14, 29, 52)
BLUE = RGBColor(34, 111, 184)
CYAN = RGBColor(47, 188, 206)
ORANGE = RGBColor(241, 137, 45)
RED = RGBColor(206, 67, 67)
GREEN = RGBColor(47, 150, 105)
LIGHT = RGBColor(241, 245, 249)
MID = RGBColor(203, 213, 225)
DARK = RGBColor(30, 41, 59)
WHITE = RGBColor(255, 255, 255)


def add_text(slide, text, x, y, w, h, size=18, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, fill=WHITE, line=MID, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def title(slide, heading, kicker=None, number=None):
    add_text(slide, heading, 0.65, 0.35, 11.8, 0.55, 27, NAVY, True)
    if kicker:
        add_text(slide, kicker, 0.68, 0.93, 11.6, 0.35, 11, BLUE, True)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.22),
                           Inches(12.0), Inches(0.035)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = CYAN
    slide.shapes[-1].line.fill.background()
    if number is not None:
        add_text(slide, f"{number:02d}", 12.25, 7.05, 0.45, 0.2, 9, MID, True,
                 PP_ALIGN.RIGHT)


def bullets(slide, items, x, y, w, h, size=17, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.font.name = "Aptos"; p.font.size = Pt(size)
        p.font.color.rgb = color; p.space_after = Pt(10)
        p.text = "•  " + p.text
    return box


def card(slide, heading, body, x, y, w, h, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, MID)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                           Inches(0.08), Inches(h)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = accent; slide.shapes[-1].line.fill.background()
    add_text(slide, heading, x + 0.25, y + 0.18, w - 0.4, 0.35, 16, NAVY, True)
    add_text(slide, body, x + 0.25, y + 0.67, w - 0.45, h - 0.82, 12.5, DARK)


def code_image(source: Path, start: int, end: int, out: Path, language: str):
    lines = source.read_text().splitlines()
    code = "\n".join(f"{i:>4}  {lines[i-1]}" for i in range(start, min(end, len(lines)) + 1))
    lexer = FortranLexer() if language == "fortran" else PythonLexer()
    formatter = ImageFormatter(
        font_name="DejaVu Sans Mono", font_size=18, line_numbers=False,
        style="monokai", image_pad=20, line_pad=4,
    )
    out.write_bytes(highlight(code, lexer, formatter))


def add_code(slide, path: Path, x, y, w, h, label):
    rect(slide, x, y, w, h, NAVY, NAVY, False)
    add_text(slide, label, x + 0.2, y + 0.12, w - 0.4, 0.25, 10, CYAN, True,
             font="DejaVu Sans Mono")
    slide.shapes.add_picture(str(path), Inches(x + 0.1), Inches(y + 0.48),
                             width=Inches(w - 0.2), height=Inches(h - 0.58))


def flow(slide, labels, y, colors=None, x0=0.65, total_width=12.0):
    colors = colors or [BLUE] * len(labels)
    n = len(labels); gap = 0.35
    width = (total_width - gap * (n - 1)) / n
    for i, (label, color) in enumerate(zip(labels, colors)):
        x = x0 + i * (width + gap)
        rect(slide, x, y, width, 0.85, color, color)
        add_text(slide, label, x + 0.08, y + 0.06, width - 0.16, 0.72,
                 13, WHITE, True, PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_text(slide, "→", x + width, y + 0.2, gap, 0.4, 20, ORANGE, True,
                     PP_ALIGN.CENTER)


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    with TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        snippets = {
            "case407": (ROOT / "cases/haas_sturtevant/generated/helium_cylinder_ms122/407.nml", 1, 24, "fortran"),
            "schlieren": (ROOT / "build/UCNS3D/src/io.f90", 57, 76, "fortran"),
            "metrics": (ROOT / "analysis/compute_ucns3d_metrics.py", 329, 353, "python"),
            "objectives": (ROOT / "optimisation/example_driver.py", 99, 112, "python"),
        }
        images = {}
        for name, (src, start, end, lang) in snippets.items():
            images[name] = tmp / f"{name}.png"
            code_image(src, start, end, images[name], lang)

        # 1 — title
        s = prs.slides.add_slide(blank); rect(s, 0, 0, 13.333, 7.5, NAVY, NAVY, False)
        add_text(s, "CFD9 + UCNS3D", 0.8, 1.25, 11.7, 0.7, 38, WHITE, True)
        add_text(s, "From a fragile solver build to a reproducible shock–bubble optimisation workflow",
                 0.82, 2.08, 10.9, 0.95, 24, CYAN, False)
        flow(s, ["BUILD", "SOLVER", "CASES", "ANALYSIS", "OPTIMISE"], 4.25,
             [BLUE, RED, ORANGE, GREEN, CYAN])
        add_text(s, "Project work after inherited CFD8 baseline • June–July 2026",
                 0.82, 6.72, 11.6, 0.3, 12, MID)

        # 2 — proposed research question
        s = prs.slides.add_slide(blank); title(s, "The proposed IRP", "ORIGINAL RESEARCH SCOPE", 2)
        rect(s, 0.85, 1.55, 11.65, 1.75, NAVY, NAVY)
        add_text(s, "Can the configuration of 2D gas cylinders be optimised to enhance shock-induced pressure amplification?",
                 1.2, 1.88, 10.95, 0.95, 25, WHITE, True, PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        card(s, "Physics objective", "Study pressure amplification from shock–bubble interactions, including strong/weak shocks, Atwood number, bubble separation, and bubble size.", 0.8, 3.75, 3.75, 2.15, RED)
        card(s, "Numerical objective", "Assess which UCNS3D capabilities are required to resolve material interfaces and the associated dynamics accurately.", 4.8, 3.75, 3.75, 2.15, BLUE)
        card(s, "Optimisation objective", "Explore multi-objective trade-offs between maximised pressure amplification and minimised material mixing.", 8.8, 3.75, 3.75, 2.15, GREEN)
        add_text(s, "Proposal: Optimisation of 2D Shock-Bubble Interaction for Pressure Amplification with UCNS3D",
                 1.0, 6.55, 11.3, 0.35, 11, MID, False, PP_ALIGN.CENTER)

        # 3 — proposed work packages
        s = prs.slides.add_slide(blank); title(s, "Proposed route from validation to optimisation", "IRP WORK PACKAGES", 3)
        flow(s, ["literature", "single-bubble convergence", "empirical validation", "new metrics", "two-bubble optimisation"], 1.55,
             [DARK, BLUE, CYAN, ORANGE, GREEN])
        card(s, "Core work", "Literature review; improved single-bubble mesh convergence; comparison with cylindrical experiments such as Haas & Sturtevant; interface, vorticity, and enstrophy metrics.", 0.75, 3.05, 5.65, 2.55, BLUE)
        card(s, "Optional extensions", "Synthetic schlieren; multi-objective pressure-versus-mixing study; 2D/3D comparison; comparison with another 2D hydrodynamics code.", 6.9, 3.05, 5.65, 2.55, ORANGE)
        add_text(s, "The implementation work in the following slides builds the infrastructure needed to execute this scope reproducibly.",
                 1.0, 6.25, 11.3, 0.55, 17, NAVY, True, PP_ALIGN.CENTER)

        # 4 — project scope
        s = prs.slides.add_slide(blank); title(s, "What changed in CFD9", "PROJECT SCOPE", 4)
        card(s, "UCNS3D reliability", "Dependency builds, MPI launch, single-rank partitioning, and VTU/PVTU output.", 0.7, 1.6, 3.8, 2.0, RED)
        card(s, "UCNS3D capability", "First-pass N-material handling, configurable bubble case 407, and synthetic schlieren.", 4.78, 1.6, 3.8, 2.0, ORANGE)
        card(s, "Python workflow", "Meshing, case generation, monitoring, convergence analysis, and pymoo/Slurm optimisation.", 8.86, 1.6, 3.8, 2.0, GREEN)
        add_text(s, "Inherited CFD8 analysis code is intentionally excluded from this presentation.",
                 1.3, 4.55, 10.7, 0.55, 20, NAVY, True, PP_ALIGN.CENTER)
        flow(s, ["outer cfd9 repo", "build/UCNS3D fork", "upstream dependency checkouts"], 5.55,
             [GREEN, BLUE, DARK])

        # 3 — build chain
        s = prs.slides.add_slide(blank); title(s, "Making the solver reproducible", "BUILD + SMOKE TEST", 5)
        flow(s, ["GKlib", "METIS", "ParMETIS", "UCNS3D", "perf2 smoke"], 1.65,
             [DARK, DARK, DARK, BLUE, GREEN])
        bullets(s, [
            "Consistent compiler and MPI toolchain across every static library.",
            "Dependency headers and archives staged without relying on fragile full installs.",
            "Fresh-directory builds verified, including a complete run under /tmp.",
            "Smoke cases retained, MPI output captured, and Git LFS pointers detected.",
            "Success now means a real timestep and valid solver output—not merely a linked binary.",
        ], 0.9, 3.0, 7.2, 3.4, 17)
        card(s, "Why this mattered", "Most early failures looked like UCNS3D defects but came from mismatched MPI, ParMETIS, METIS, or GKlib builds.", 8.55, 3.15, 3.8, 2.2, RED)

        # 4 — bugs
        s = prs.slides.add_slide(blank); title(s, "The failure modes we chased", "UCNS3D DEBUGGING", 6)
        card(s, "MPI launch", "PRRTE executable lookup, sandbox interfaces, and misleading missing help files.", 0.7, 1.55, 3.8, 1.45, RED)
        card(s, "Partitioning", "-np 1 entered distributed assumptions and indexed partition arrays at zero.", 4.78, 1.55, 3.8, 1.45, ORANGE)
        card(s, "Parallel output", "Truncated appended VTU data, temporary files, invalid frees, and mode-5 regressions.", 8.86, 1.55, 3.8, 1.45, RED)
        card(s, "Input integrity", "Git LFS pointer text was sometimes copied and treated as a real mesh.", 0.7, 3.35, 3.8, 1.45, ORANGE)
        card(s, "Platform mismatch", "Local success but production SIGABRT around ParMETIS, libgomp, or output.", 4.78, 3.35, 3.8, 1.45, RED)
        card(s, "Test quality", "File existence was insufficient; tests now require time advancement and coherent outputs.", 8.86, 3.35, 3.8, 1.45, GREEN)
        add_text(s, "Outcome: both -np 1 and -np 2 reached time stepping; production-only invalid frees still need final root-cause isolation.",
                 0.9, 5.65, 11.6, 0.8, 18, NAVY, True, PP_ALIGN.CENTER)

        # 5 — N material
        s = prs.slides.add_slide(blank); title(s, "From two-material assumptions toward N materials", "ALLAIRE MODEL", 7)
        add_text(s, "Original state", 0.75, 1.5, 2.5, 0.4, 18, RED, True)
        bullets(s, ["N-sized arrays", "Some N-species EOS loops", "Hard-coded source index", "Two-material boundaries", "Two-field output assumptions"], 0.8, 2.0, 3.0, 3.6, 14)
        add_text(s, "→", 3.7, 3.0, 0.8, 0.6, 34, ORANGE, True, PP_ALIGN.CENTER)
        add_text(s, "First-pass generalisation", 4.55, 1.5, 3.5, 0.4, 18, BLUE, True)
        bullets(s, ["Looped nonconservative terms", "N-aware profiles", "Three-material test", "Reconstruct final αₙ", "Write all N volume fractions"], 4.6, 2.0, 3.3, 3.6, 14)
        rect(s, 8.5, 1.45, 4.0, 4.55, LIGHT, MID)
        add_text(s, "State concept", 8.8, 1.75, 3.4, 0.35, 17, NAVY, True, PP_ALIGN.CENTER)
        add_text(s, "αρ₁ … αρₙ", 9.15, 2.35, 2.7, 0.55, 25, BLUE, True, PP_ALIGN.CENTER)
        add_text(s, "+", 9.15, 2.95, 2.7, 0.4, 22, DARK, True, PP_ALIGN.CENTER)
        add_text(s, "α₁ … αₙ₋₁", 9.15, 3.4, 2.7, 0.55, 25, ORANGE, True, PP_ALIGN.CENTER)
        add_text(s, "αₙ = 1 − Σ αₖ", 9.15, 4.35, 2.7, 0.55, 18, GREEN, True, PP_ALIGN.CENTER)
        add_text(s, "Functional first pass—not yet formal numerical verification.", 8.85, 5.15, 3.3, 0.55, 12, RED, True, PP_ALIGN.CENTER)

        # 6 — 407
        s = prs.slides.add_slide(blank); title(s, "Case 407: bubbles become data, not source code", "CONFIGURABLE INITIAL CONDITION", 8)
        add_code(s, images["case407"], 0.65, 1.48, 7.35, 5.45, "cases/haas_sturtevant/generated/helium_cylinder_ms122/407.nml")
        card(s, "407.nml", "Arbitrary bubble count; centre, radius, material state, velocity, pressure, and perturbation controls.", 8.35, 1.55, 4.25, 1.75, BLUE)
        card(s, "MULTISPECIES.DAT", "Defines the EOS and thermodynamic properties of the materials referenced by the profile.", 8.35, 3.55, 4.25, 1.45, ORANGE)
        card(s, "UCNS3D.DAT", "Retains solver, numerics, time-stepping, and output configuration.", 8.35, 5.25, 4.25, 1.35, GREEN)

        # 7 — schlieren
        s = prs.slides.add_slide(blank); title(s, "Synthetic schlieren added to solver output", "DENSITY-GRADIENT DIAGNOSTIC", 9)
        add_code(s, images["schlieren"], 0.65, 1.48, 7.0, 5.25, "build/UCNS3D/src/io.f90")
        rect(s, 8.0, 1.6, 4.55, 2.0, LIGHT, MID)
        add_text(s, "S = |∇ρ|", 8.4, 2.05, 3.75, 0.65, 32, BLUE, True, PP_ALIGN.CENTER)
        add_text(s, "Computed per cell and written with VTU solution fields", 8.45, 2.85, 3.65, 0.45, 14, DARK, False, PP_ALIGN.CENTER)
        bullets(s, ["Direct ParaView/VisIt access", "Inverted X-ray VisIt view", "Comparable framing for experiment imagery"], 8.2, 4.15, 4.1, 2.1, 15)

        # 8 — python generation
        s = prs.slides.add_slide(blank); title(s, "Python makes complete experiments reproducible", "MESH + CASE GENERATION", 10)
        flow(s, ["physical setup", "mesh spacing", "unstructured mesh", "UCNS3D inputs", "Slurm case"], 1.55,
             [ORANGE, CYAN, BLUE, GREEN, DARK])
        card(s, "Meshing", "Open-source unstructured channel meshes with local interaction-region refinement and coarser inlet/outlet zones.", 0.75, 3.05, 3.75, 2.25, BLUE)
        card(s, "Generic cases", "N bubbles, material selection, optional RTP density, shock state, domain size, and physical mesh spacing.", 4.8, 3.05, 3.75, 2.25, GREEN)
        card(s, "Reference studies", "Haas–Sturtevant helium/R22 wrappers plus a literature-based water–air convergence case.", 8.85, 3.05, 3.75, 2.25, ORANGE)
        add_text(s, "One Python driver can regenerate the mesh, inputs, scheduler file, and case documentation.",
                 1.0, 6.0, 11.3, 0.5, 19, NAVY, True, PP_ALIGN.CENTER)

        # 9 — monitor/convergence
        s = prs.slides.add_slide(blank); title(s, "Runs become observable and comparable", "MONITORING + CONVERGENCE", 11)
        card(s, "Recursive discovery", "Find run directories from UCNS3D.DAT and mesh/output evidence.", 0.7, 1.55, 3.75, 1.55, BLUE)
        card(s, "Portable status", "Classify not started, running, failed, or successful—with or without Slurm.", 0.7, 3.35, 3.75, 1.55, GREEN)
        card(s, "Physical convergence", "Process mesh_* suites and report converged physical mesh spacing, not just total cells.", 0.7, 5.15, 3.75, 1.55, ORANGE)
        flow(s, ["generate", "submit", "monitor", "analyse", "reuse Δx"], 2.0,
             [BLUE, DARK, GREEN, ORANGE, CYAN], x0=4.9, total_width=7.65)
        add_text(s, "Same resolution criterion can be transferred to a different domain extent.",
                 5.15, 4.4, 6.8, 0.7, 21, NAVY, True, PP_ALIGN.CENTER)

        # 10 — optimisation
        s = prs.slides.add_slide(blank); title(s, "Multi-objective optimisation closes the loop", "PYMOO + SLURM", 12)
        flow(s, ["ask", "generate cases", "submit jobs", "read metrics", "tell"], 1.55,
             [CYAN, BLUE, DARK, GREEN, ORANGE])
        add_code(s, images["objectives"], 0.75, 3.0, 6.7, 3.6, "optimisation/example_driver.py")
        bullets(s, ["Python-native parameter bounds", "One Slurm job per population point", "No Dask dependency", "Python 3.12 via uv", "Maximisation mapped to pymoo minimisation signs"], 8.0, 3.15, 4.5, 3.15, 16)

        # 11 — metrics
        s = prs.slides.add_slide(blank); title(s, "Robust objectives from VTU/PVTU time series", "OPTIMISATION METRICS", 13)
        add_code(s, images["metrics"], 0.65, 1.45, 7.15, 5.55, "analysis/compute_ucns3d_metrics.py")
        card(s, "Ap95_target ↑", "Peak-in-time target 95th-percentile pressure amplification; avoids one-cell spikes.", 8.15, 1.55, 4.4, 1.15, BLUE)
        card(s, "Ip_target ↑", "Area-weighted target excess-pressure impulse integrated over time.", 8.15, 2.95, 4.4, 1.15, GREEN)
        card(s, "Mbad_target ↓", "Maximum unwanted-material fraction in the desired compression region.", 8.15, 4.35, 4.4, 1.15, RED)
        card(s, "Lp_localisation ↑", "Optional ratio of target excess pressure to the larger interaction region.", 8.15, 5.75, 4.4, 1.15, ORANGE)

        # 12 — status
        s = prs.slides.add_slide(blank); title(s, "Where the project stands", "OUTCOMES + NEXT VALIDATION", 14)
        card(s, "Working", "Reproducible build and smoke path; case 407; N-field output; synthetic schlieren; meshing/case generation; monitoring; optimisation bridge.", 0.7, 1.55, 3.8, 3.1, GREEN)
        card(s, "Needs validation", "Formal N-material Allaire verification; conservation tests; reference solutions; cross-platform VTU stress testing.", 4.78, 1.55, 3.8, 3.1, ORANGE)
        card(s, "Remaining risk", "Production-only invalid frees suggest toolchain/runtime mismatch or latent parallel-I/O corruption.", 8.86, 1.55, 3.8, 3.1, RED)
        add_text(s, "The central achievement", 1.0, 5.35, 11.3, 0.4, 17, BLUE, True, PP_ALIGN.CENTER)
        add_text(s, "A repeatable path from physical case definition → UCNS3D simulation → quantitative optimisation objectives",
                 1.0, 5.9, 11.3, 0.85, 23, NAVY, True, PP_ALIGN.CENTER)

        OUT.parent.mkdir(parents=True, exist_ok=True)
        prs.save(OUT)


if __name__ == "__main__":
    build()
